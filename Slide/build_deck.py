"""Builds the P01 pitch deck.

Aesthetic is taken from Slide/ReferenceSlide.pptx: a warm cream ground, a single
amber accent, and heavy black display type. Nothing is copied from that file --
the palette and proportions are rebuilt here with drawn shapes and our own
screenshots.
"""

import io
import json
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")

# ----------------------------------------------------------------- palette
CREAM = RGBColor(0xFF, 0xFD, 0xF8)
AMBER = RGBColor(0xF8, 0xC2, 0x38)
INK = RGBColor(0x11, 0x11, 0x11)
INK_SOFT = RGBColor(0x5D, 0x68, 0x80)
# Amber reads at only ~1.9:1 on cream, so highlighted words use a deeper bronze.
AMBER_TEXT = RGBColor(0x8A, 0x5D, 0x05)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE8, 0xE3, 0xD6)
CUT = RGBColor(0xEF, 0x6B, 0x5E)
GRID = RGBColor(0x4C, 0x6F, 0xFF)
GEN = RGBColor(0x1F, 0xAE, 0x87)

CREAM_HEX = "#FFFDF8"
AMBER_HEX = "#F8C238"
INK_HEX = "#111111"

W, H = Inches(20), Inches(11.25)
M = Inches(1.35)          # page margin
CONTENT_W = W - 2 * M

for weight in ("Regular", "Medium", "SemiBold", "Bold", "ExtraBold"):
    path = os.path.join(ASSETS, f"Poppins-{weight}.ttf")
    if os.path.exists(path):
        fm.fontManager.addfont(path)

stats = json.load(io.open(os.path.join(ASSETS, "stats.json"), encoding="utf-8"))
ROWS, TOTALS = stats["rows"], stats["totals"]
AVG_CUTS = sum(r["cuts"] for r in ROWS) / len(ROWS)
AVG_CUT_MIN = sum(r["cutMinutes"] for r in ROWS) / len(ROWS)


# ------------------------------------------------------------------- charts
def chart_placed_donut(path):
    fig, ax = plt.subplots(figsize=(6.2, 6.2), dpi=200)
    fig.patch.set_facecolor(CREAM_HEX)
    ax.set_facecolor(CREAM_HEX)
    placed, impossible = TOTALS["placed"], TOTALS["unplaced"]
    wedges, _ = ax.pie(
        [placed, impossible],
        startangle=90,
        colors=[AMBER_HEX, "#E4DFD2"],
        wedgeprops=dict(width=0.30, edgecolor=CREAM_HEX, linewidth=6),
    )
    ax.text(0, 0.12, f"{placed}", ha="center", va="center",
            fontsize=58, fontweight="bold", family="Poppins", color=INK_HEX)
    ax.text(0, -0.20, "of 263 jobs placed", ha="center", va="center",
            fontsize=15, family="Poppins", color="#5D6880")
    ax.set(aspect="equal")
    fig.savefig(path, transparent=False, bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)


def chart_generator_bars(path):
    fig, ax = plt.subplots(figsize=(15.5, 3.9), dpi=200)
    fig.patch.set_facecolor(CREAM_HEX)
    ax.set_facecolor(CREAM_HEX)
    labels = [r["caseId"].replace("PUB-", "") for r in ROWS]
    values = [r["generator"] for r in ROWS]
    peak = max(values)
    colors = [AMBER_HEX if v == peak else "#EBDFC0" for v in values]
    ax.bar(labels, values, color=colors, edgecolor="none", width=0.68)
    ax.set_ylabel("generator minutes", family="Poppins", fontsize=13, color="#5D6880")
    ax.tick_params(axis="both", labelsize=11, colors="#5D6880", length=0)
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color("#DDD6C6")
    ax.grid(axis="y", color="#E8E3D6", linewidth=1)
    ax.set_axisbelow(True)
    for label in ax.get_xticklabels() + ax.get_yticklabels():
        label.set_family("Poppins")
    fig.tight_layout()
    fig.savefig(path, facecolor=CREAM_HEX, bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)


# ------------------------------------------------------------ device frames
def rounded(img, radius):
    mask = Image.new("L", img.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle([0, 0, img.width - 1, img.height - 1], radius, fill=255)
    out = Image.new("RGBA", img.size, (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)
    return out


def laptop_frame(shot_path, out_path, screen_w=2400):
    shot = Image.open(shot_path).convert("RGB")
    ratio = 10 / 16
    target_h = int(screen_w * ratio)
    # Crop the top of the page to a 16:10 window rather than squashing it.
    crop_h = min(shot.height, int(shot.width * ratio))
    shot = shot.crop((0, 0, shot.width, crop_h)).resize((screen_w, target_h), Image.LANCZOS)

    bezel = 26
    body_w, body_h = screen_w + bezel * 2, target_h + bezel * 2
    base_h, base_extra = 46, 150
    canvas = Image.new("RGBA", (body_w + base_extra * 2, body_h + base_h + 30), (0, 0, 0, 0))
    draw = ImageDraw.Draw(canvas)

    draw.rounded_rectangle(
        [base_extra, 0, base_extra + body_w, body_h], 34, fill=(17, 17, 17, 255)
    )
    canvas.paste(rounded(shot, 12), (base_extra + bezel, bezel), rounded(shot, 12))
    draw.rounded_rectangle(
        [0, body_h + 6, canvas.width - 1, body_h + base_h], 18, fill=(30, 30, 30, 255)
    )
    draw.rounded_rectangle(
        [canvas.width // 2 - 110, body_h + 6, canvas.width // 2 + 110, body_h + 22],
        8, fill=(58, 58, 58, 255),
    )
    canvas.save(out_path)


def phone_frame(shot_path, out_path, screen_w=760):
    shot = Image.open(shot_path).convert("RGB")
    target_h = int(screen_w * (shot.height / shot.width))
    shot = shot.resize((screen_w, target_h), Image.LANCZOS)

    bezel = 22
    canvas = Image.new("RGBA", (screen_w + bezel * 2, target_h + bezel * 2), (0, 0, 0, 0))
    draw = ImageDraw.Draw(canvas)
    draw.rounded_rectangle([0, 0, canvas.width - 1, canvas.height - 1], 74, fill=(17, 17, 17, 255))
    canvas.paste(rounded(shot, 54), (bezel, bezel), rounded(shot, 54))
    cx = canvas.width // 2
    draw.rounded_rectangle([cx - 78, bezel + 14, cx + 78, bezel + 44], 15, fill=(17, 17, 17, 255))
    canvas.save(out_path)


def place_picture(slide, name, x, y, width):
    """Adds a picture at its true aspect ratio and returns its bottom edge."""
    path = os.path.join(ASSETS, name)
    with Image.open(path) as im:
        aspect = im.height / im.width
    height = Emu(int(width * aspect))
    slide.shapes.add_picture(path, x, y, width=width, height=height)
    return y + height


def asset(slide, name, x, y, size):
    """Drops a square 3D element from the reference set."""
    slide.shapes.add_picture(os.path.join(ASSETS, name), x, y, width=size, height=size)


# --------------------------------------------------------------- pptx tools
def add_slide(prs, bg=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def box(slide, x, y, w, h, text_parts, size=20, align=PP_ALIGN.LEFT,
        color=INK, font="Poppins", bold=False, line=1.25, space_after=0):
    """`text_parts` is a string, or a list of (text, style) where style is one
    of '', 'b' (bold), 'a' (amber), 'm' (muted)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paragraphs = text_parts if isinstance(text_parts, list) else [[(text_parts, "")]]
    if paragraphs and isinstance(paragraphs[0], tuple):
        paragraphs = [paragraphs]

    for index, parts in enumerate(paragraphs):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        para.space_after = Pt(space_after)
        for text, style in parts:
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.name = font
            run.font.bold = bold or "b" in style
            if "a" in style:
                run.font.color.rgb = AMBER_TEXT
            elif "m" in style:
                run.font.color.rgb = INK_SOFT
            else:
                run.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, fill=CARD, radius=True, line_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def eyebrow(slide, text, y=None, color=INK_SOFT):
    box(slide, M, y or Inches(0.95), CONTENT_W, Inches(0.4),
        [[(text.upper(), "b")]], size=15, color=color)


def heading(slide, text, y=Inches(1.45), size=58, color=INK):
    box(slide, M, y, CONTENT_W, Inches(1.4), [[(text, "b")]], size=size, color=color, line=1.05)


def rule(slide, y, w=Inches(2.6), color=AMBER):
    bar = rect(slide, M, y, w, Inches(0.13), fill=color, radius=False)
    return bar


def stat_card(slide, x, y, w, h, value, label, accent=False):
    rect(slide, x, y, w, h, fill=AMBER if accent else CARD,
         line_color=None if accent else LINE)
    box(slide, x + Inches(0.42), y + Inches(0.38), w - Inches(0.84), Inches(1.0),
        [[(value, "b")]], size=44, color=INK, line=1.0)
    box(slide, x + Inches(0.42), y + h - Inches(0.95), w - Inches(0.84), Inches(0.8),
        [[(label, "")]], size=15, color=INK if accent else INK_SOFT, line=1.2)


def numbered_card(slide, x, y, w, h, badge, title, body_parts):
    rect(slide, x, y, w, h, fill=CARD, line_color=LINE)
    asset(slide, badge, x + Inches(0.42), y + Inches(0.42), Inches(0.78))
    box(slide, x + Inches(0.42), y + Inches(1.45), w - Inches(0.84), Inches(1.1),
        [[(title, "b")]], size=20, color=INK, line=1.15)
    box(slide, x + Inches(0.42), y + Inches(2.5), w - Inches(0.84), h - Inches(2.95),
        body_parts, size=16, color=INK_SOFT, line=1.35)


# ------------------------------------------------------------------- slides
def build():
    chart_placed_donut(os.path.join(ASSETS, "chart-donut.png"))
    chart_generator_bars(os.path.join(ASSETS, "chart-generator.png"))
    laptop_frame(os.path.join(ASSETS, "shot-desktop.png"), os.path.join(ASSETS, "frame-laptop.png"))
    phone_frame(os.path.join(ASSETS, "shot-phone.png"), os.path.join(ASSETS, "frame-phone.png"))

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 -------------------------------------------------------------- title
    s = add_slide(prs, AMBER)
    box(s, M, Inches(3.1), CONTENT_W, Inches(0.5), [[("LOFISTACK HACKATHON 2026   ·   PROBLEM P01", "b")]], size=17)
    box(s, M, Inches(3.8), Inches(15.5), Inches(3.2),
        [[("Load-Shedding", "b")], [("Window Planner", "b")]], size=104, line=1.02)
    rect(s, M, Inches(7.6), Inches(3.0), Inches(0.16), fill=INK, radius=False)
    box(s, M, Inches(8.1), Inches(13), Inches(1.4),
        [[("Team Logarithm", "b"), ("   ·   LSH26-T033", "")],
         [("Oitijya Islam Auvro  ·  Md. Nafiz Ahmed  ·  Dewan Salman Rahman Zisan", "")]],
        size=20, line=1.5)

    # 2 ------------------------------------------------------------ problem
    s = add_slide(prs)
    eyebrow(s, "01  ·  The problem")
    heading(s, "Load shedding wastes the working day.")
    rule(s, Inches(2.9))
    box(s, M, Inches(3.4), Inches(9.6), Inches(3.4),
        [[("A print shop runs on machines that need ", ""), ("mains power", "ab"),
          (". When the grid goes down mid-job, the work is ", ""), ("lost", "b"),
          (" — paper wasted, deadline missed, customer waiting.", "")],
         [("Cuts are announced in advance. The information is there. What is missing is a way to ", ""),
          ("plan the day around them", "ab"), (".", "")]],
        size=22, color=INK_SOFT, line=1.45, space_after=18)

    x = M + Inches(10.6)
    stat_card(s, x, Inches(3.4), Inches(3.6), Inches(3.0), f"{AVG_CUTS:.1f}", "power cuts on an average day", accent=True)
    stat_card(s, x + Inches(3.95), Inches(3.4), Inches(3.6), Inches(3.0),
              f"{AVG_CUT_MIN / 60:.1f}h", "of the day with no grid power")
    stat_card(s, x, Inches(6.7), Inches(7.55), Inches(2.6), "263", "jobs across the 25 published cases we tested against")

    # 3 ------------------------------------------------- the required items
    s = add_slide(prs)
    eyebrow(s, "02  ·  What the problem asks for")
    heading(s, "Four required items.")
    rule(s, Inches(2.9))
    items = [
        ("badge-1.png", "Enter the cuts", [[("Start and end time for each announced cut, drawn on a ", ""), ("24-hour timeline", "ab"), (".", "")]]),
        ("badge-2.png", "Add the jobs", [[("Name, duration, and what power it needs: ", ""), ("grid", "ab"), (", ", ""), ("generator", "ab"), (" or ", ""), ("none", "ab"), (".", "")]]),
        ("badge-3.png", "Place them automatically", [[("No job needing grid power may ", ""), ("ever", "b"), (" land inside a cut. Plan shown beside the cut bars.", "")]]),
        ("badge-4.png", "Count generator minutes", [[("The total the plan needs, updating the ", ""), ("moment", "ab"), (" a job is added or removed.", "")]]),
    ]
    cw = (CONTENT_W - Inches(0.75)) / 4
    for i, (badge, title, body) in enumerate(items):
        numbered_card(s, M + i * (cw + Inches(0.25)), Inches(3.4), cw, Inches(5.9), badge, title, body)

    # 4 ----------------------------------------------------------- solution
    s = add_slide(prs)
    eyebrow(s, "03  ·  Our solution")
    heading(s, "One screen. The whole day.")
    s.shapes.add_picture(os.path.join(ASSETS, "frame-laptop.png"),
                         M - Inches(0.3), Inches(3.2), width=Inches(11.6))
    bx = M + Inches(11.9)
    box(s, bx, Inches(3.3), Inches(5.4), Inches(5.5),
        [[("Cuts on top. The plan directly beneath, on the ", ""), ("same scale", "ab"),
          (" — so a judge can see at a glance that no blue grid job touches a red cut.", "")],
         [("Amber marks the minutes a job actually spends ", ""), ("on the generator", "ab"), (".", "")],
         [("Every placement carries a one-line reason in plain English.", "")],
         [("No login. No database. No API key.", "b")]],
        size=19, color=INK_SOFT, line=1.4, space_after=16)

    # 5 ---------------------------------------------------------- algorithm
    s = add_slide(prs)
    eyebrow(s, "04  ·  How the plan is built")
    heading(s, "Most-constrained first.")
    rule(s, Inches(2.9))
    steps = [
        ("badge-a1.png", "Grid jobs", [[("Cut-free space only. A ", ""), ("hard rule", "ab"),
          (" — if no gap is long enough the job is reported, never placed illegally.", "")]]),
        ("badge-a2.png", "Generator jobs", [[("The slot costing the ", ""), ("fewest generator minutes", "ab"),
          (", not simply the earliest — because that cost is exactly what R4 reports.", "")]]),
        ("badge-a3.png", "No-power jobs", [[("Parked ", ""), ("inside", "b"), (" the cuts on purpose. A cut costs them nothing, which keeps scarce cut-free time for the rest.", "")]]),
    ]
    cw = (CONTENT_W - Inches(0.6)) / 3
    for i, (badge, title, body) in enumerate(steps):
        numbered_card(s, M + i * (cw + Inches(0.3)), Inches(3.4), cw, Inches(5.2), badge, title, body)
    box(s, M, Inches(9.05), CONTENT_W, Inches(1.6),
        [[("Within each group: ", ""), ("rush orders first", "ab"),
          (", then ", ""), ("earliest promised time", "ab"),
          (", then longest. The result is ", ""), ("deterministic", "b"),
          (" — the same input always gives the same plan.", "")]],
        size=20, color=INK_SOFT, line=1.4)

    # 6 ------------------------------------------------------- architecture
    s = add_slide(prs)
    eyebrow(s, "05  ·  Architecture")
    heading(s, "The rules are not in the interface.")
    rule(s, Inches(2.9))
    layers = [
        ("src/domain/", "Pure TypeScript. No React.",
         [[("time · intervals · schedule · fixture", "b")],
          [("Every rule lives here, so it can be tested directly and read without running anything.", "")]], AMBER),
        ("src/components/", "React, presentation only.",
         [[("Timeline · Inputs · PlanPanel", "b")],
          [("Holds cuts, jobs and hours. Nothing else.", "")]], CARD),
        ("The plan", "Derived, never stored.",
         [[("buildPlan(state) on every render", "b")],
          [("There is no second copy of the generator total to go stale — R4 cannot drift.", "")]], CARD),
    ]
    cw = (CONTENT_W - Inches(0.6)) / 3
    for i, (title, sub, body, fill) in enumerate(layers):
        x = M + i * (cw + Inches(0.3))
        rect(s, x, Inches(3.45), cw, Inches(4.3), fill=fill, line_color=None if fill == AMBER else LINE)
        box(s, x + Inches(0.5), Inches(3.85), cw - Inches(1.0), Inches(0.7), [[(title, "b")]], size=24)
        box(s, x + Inches(0.5), Inches(4.6), cw - Inches(1.0), Inches(0.6), [[(sub, "b")]], size=16,
            color=INK if fill == AMBER else INK_SOFT)
        box(s, x + Inches(0.5), Inches(5.4), cw - Inches(1.0), Inches(2.0), body, size=15,
            color=INK if fill == AMBER else INK_SOFT, line=1.35, space_after=10)

    # 7 ------------------------------------------------------- real shop
    s = add_slide(prs)
    eyebrow(s, "06  ·  Beyond the four items")
    heading(s, "Built for a real shop.")
    rule(s, Inches(2.9))
    extras = [
        ("Machines", [[("The shop runs ", ""), ("more than one", "ab"), (" job at a time. Set 1 to 8 lanes; every lane gets its own row on the timeline.", "")]]),
        ("Promised times", [[("A customer collects at 5pm. A job that ", ""), ("cannot", "b"), (" make its promise is reported, never quietly placed late.", "")]]),
        ("Rush orders", [[("Mark a job urgent and it takes the ", ""), ("front of the day", "ab"), (", ahead of everything else in its class.", "")]]),
        ("Fuel budget", [[("Diesel is bought by the litre. Cap the generator minutes and the planner ", ""), ("will not exceed it", "ab"), (".", "")]]),
    ]
    cw = (CONTENT_W - Inches(0.75)) / 4
    for i, (title, body) in enumerate(extras):
        x = M + i * (cw + Inches(0.25))
        rect(s, x, Inches(3.45), cw, Inches(5.85), fill=CARD, line_color=LINE)
        rect(s, x + Inches(0.45), Inches(3.95), Inches(0.9), Inches(0.14), fill=AMBER, radius=False)
        box(s, x + Inches(0.45), Inches(4.45), cw - Inches(0.9), Inches(0.9), [[(title, "b")]], size=23)
        box(s, x + Inches(0.45), Inches(5.6), cw - Inches(0.9), Inches(3.4), body, size=16,
            color=INK_SOFT, line=1.35)

    # 8 ----------------------------------------------------------- testing
    s = add_slide(prs)
    eyebrow(s, "07  ·  Proof")
    heading(s, "Proven, not claimed.")
    rule(s, Inches(2.9))
    asset(s, "mark-tick.png", W - M - Inches(1.5), Inches(9.45), Inches(1.5))
    box(s, M, Inches(3.4), Inches(9.4), Inches(5.5),
        [[("69 automated tests", "ab")],
         [("Four rules are asserted for every plan, on hand-written cases and on all ", ""), ("25 published cases", "b"), (":", "")],
         [("· a grid job never overlaps a cut, by even one minute", "")],
         [("· no two jobs overlap on the same machine", "")],
         [("· generator minutes equal the time actually spent inside a cut", "")],
         [("· every placement respects its ready and promised times", "")],
         [("Writing the tests first found a real defect: the planner had been giving generator jobs the earliest slot instead of the cheapest, ", ""), ("inflating the R4 total", "ab"), (".", "")]],
        size=19, color=INK_SOFT, line=1.4, space_after=14)
    rect(s, M + Inches(10.2), Inches(3.3), Inches(7.1), Inches(5.9), fill=CARD, line_color=LINE)
    s.shapes.add_picture(os.path.join(ASSETS, "chart-donut.png"),
                         M + Inches(11.85), Inches(3.6), height=Inches(4.2), width=Inches(4.2))
    box(s, M + Inches(10.7), Inches(8.05), Inches(6.1), Inches(1.0),
        [[("The other 36 are genuinely impossible — jobs longer than the whole day, or needing more uninterrupted cut-free time than exists.", "")]],
        size=14, color=INK_SOFT, align=PP_ALIGN.CENTER, line=1.3)

    # 9 -------------------------------------------------------- measurement
    s = add_slide(prs)
    eyebrow(s, "08  ·  Measured, not guessed")
    heading(s, "Generator minutes, case by case.")
    bottom = place_picture(s, "chart-generator.png", M, Inches(3.15), CONTENT_W)
    box(s, M, bottom + Inches(0.45), CONTENT_W, Inches(1.4),
        [[("1,845 generator minutes", "ab"), (" across 25 cases. Every number here comes from running the real engine over the organisers' own fixture — ", ""), ("no case was hand-picked", "b"), (".", "")]],
        size=19, color=INK_SOFT, line=1.4)

    # 10 ---------------------------------------------------------- the demo
    s = add_slide(prs)
    eyebrow(s, "09  ·  Live")
    heading(s, "Try it yourself.")
    s.shapes.add_picture(os.path.join(ASSETS, "frame-phone.png"),
                         M + Inches(0.4), Inches(3.1), height=Inches(7.0))
    bx = M + Inches(6.4)
    asset(s, "mark-cursor.png", W - M - Inches(2.0), Inches(9.15), Inches(1.9))
    box(s, bx, Inches(3.4), Inches(10.5), Inches(1.0),
        [[("lsh26-t033-p01.vercel.app", "ab")]], size=36)
    box(s, bx, Inches(4.6), Inches(10.5), Inches(4.4),
        [[("github.com/AuvroIslam/lsh26-t033-p01", "b")],
         [("Opens with no account and no key. Press ", ""), ("Load published sample data", "b"),
          (" and all 25 cases load; upload any file in the same shape to try your own.", "")],
         [("Works on a phone. State is kept in the browser, and Reset restores it.", "")]],
        size=19, color=INK_SOFT, line=1.4, space_after=16)

    # 11 -------------------------------------------------------- thank you
    s = add_slide(prs, AMBER)
    box(s, M, Inches(4.3), Inches(15), Inches(2.4), [[("Thank you.", "b")]], size=104, line=1.0)
    rect(s, M, Inches(6.9), Inches(3.0), Inches(0.16), fill=INK, radius=False)
    box(s, M, Inches(7.4), Inches(14), Inches(1.6),
        [[("Team Logarithm", "b"), ("   ·   LSH26-T033   ·   P01 Load-Shedding Window Planner", "")],
         [("lsh26-t033-p01.vercel.app", "")]],
        size=20, line=1.5)

    out = os.path.join(HERE, "LSH26-T033-P01.pptx")
    tmp = os.path.join(HERE, ".build.pptx")
    prs.save(tmp)
    try:
        os.replace(tmp, out)
    except PermissionError:
        print("!! Close the deck in PowerPoint -- it is locked. Built copy left at", tmp)
        return
    print("saved", out, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")


if __name__ == "__main__":
    build()
